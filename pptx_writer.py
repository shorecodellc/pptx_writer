#kevin fink
#kevin@shorecode.org
#Mar 20 06:20:56 PM +07 2025
#pptx_writer.py

from typing import Union
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

class PptxWriter(Presentation):
    '''
    Convenience class inherited from python-pptx Presentation class
    '''
    def __init__(filepath):
        super().__init__()
    
    @staticmethod
    def change_font_tf(shape, font):
        """
        Change the font family for all runs and paragraphs in a text frame for a given shape
        """
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:            
                run.font.name = font
                
    @staticmethod
    def change_font_run(shape, font, para_idx, run_idx):
        """
        Change the font family for a specified run in a text frame for a given shape
        """
        shape.text_frame.paragraphs[para_idx][run_idx]

    @staticmethod
    def change_font_para(shape, font, para_idx):
        """
        Change the font family for all runs in a specified paragraph in a text frame for a given shape
        """
        for paragraph in shape.text_frame.paragraphs:
            paragraph[para_idx]
            
    @staticmethod
    def change_font_color_tf(shape, color: list[int, int, int]):
        """
        Change the font color for all runs and paragraphs in a text frame for a given shape
        """
        rgb = RGBColor(*color)
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:            
                run.font.color.rgb = rgb
                
    @staticmethod
    def change_font_color_run(shape, color: list[int, int, int], para_idx, run_idx):
        """
        Change the font color for a specified run in a text frame for a given shape
        """
        rgb = RGBColor(*color)
        shape.text_frame.paragraphs[para_idx].runs[run_idx].font.color.rgb = rgb
            
    @staticmethod
    def change_font_color_para(shape, color: list[int, int, int], para_idx, ):
        """
        Change the font color for all runs in a specified paragraph in a text frame for a given shape
        """
        rgb = RGBColor(*color)
        for run in shape.text_frame.paragraphs[para_idx]:   
            run.font.color.rgb = rgb
            
    @staticmethod
    def change_font_size_tf(shape, font_size: Pt):
        """
        Change the font size for all runs and paragraphs in a text frame for a given shape
        """        
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:            
                run.font.size = Pt(font_size)
                
    @staticmethod
    def change_font_size_run(shape, font_size: Pt, para_idx, run_idx):
        """
        Change the font size for a specified run in a text frame for a given shape
        """
        shape.text_frame.paragraphs[para_idx].runs[run_idx].font.size = Pt(font_size)
        
    @staticmethod
    def change_font_size_para(shape, font_size: Pt, para_idx, ):
        """
        Change the font size for all runs in a specified paragraph in a text frame for a given shape
        """
        for run in shape.text_frame.paragraphs[para_idx]:  
            run.font.size = Pt(font_size)
            
    @staticmethod
    def make_vertical_center_tf(shape):
        for paragraph in shape.text_frame.paragraphs:
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
    @staticmethod
    def make_vertical_center_tf(shape, para_idx):
        shape.paragraphs[para_idx].vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
    @staticmethod
    def make_horiz_center_tf(shape):
        """
        Change the horizontal alignment of text in a text frame for a given shape
        """
        for paragraph in shape.text_frame.paragraphs:     
            paragraph.alignment = PP_ALIGN.CENTER        
            
    @staticmethod
    def make_horiz_center_para(shape, para_idx):
        """
        Change the horizontal alignment of text in a specified paragraph for a given shape
        """
        shape.text_frame.paragraphs[para_idx].alignment = PP_ALIGN.CENTER
        
    @staticmethod
    def toggle_bold_tf(shape, bold_boolean=None):
        """
        Toggle the bold property of text in a text frame for a given shape
        If `bold_boolean` is specified, the property will be set to that value
        """
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if bold_boolean == None:                    
                    run.font.bold = run.font.bold == False
                else:
                    run.font.bold = bold_boolean
                
    @staticmethod
    def toggle_bold_run(shape, para_idx, run_idx, bold_boolean=None):
        """
        Toggle the bold property of text in a specified run for a given shape
        If `bold_boolean` is specified, the property will be set to that value
        """
        if bold_boolean == None:
            shape.text_frame.paragraphs[para_idx].runs[run_idx].font.bold = shape.text_frame.paragraphs[para_idx].runs[run_idx].font.bold == False
        else:
            shape.text_frame.paragraphs[para_idx].runs[run_idx].font.bold = bold_boolean
                    
    @staticmethod
    def toggle_bold_para(shape, para_idx, bold_boolean=None):
        """
        Toggle the bold property of text for all runs in a specified paragraph for a given shape
        If `bold_boolean` is specified, the property will be set to that value
        """        
        for run in shape.text_frame.paragraphs[para_idx].runs:
            if bold_boolean == None:                    
                run.font.bold = run.font.bold == False
            else:
                run.font.bold = bold_boolean
       
    @staticmethod         
    def convert_to_string(entry):
        '''
        Used to convert malformed chatGPT responses to string
        Response should be queried to be in format {"key": "desired string data"}
        '''
        if isinstance(entry, str):
            return entry
        if isinstance(entry, dict):
            result = str(entry.values())[0]
        if isinstance(entry, list):
            for i, e in enumerate(entry):
                if isinstance(e, dict):
                    entry[i] = list(e.values())[0]        
            result = ' '.join(entry)
        return self.convert_to_string(result)
    
    @staticmethod
    def add_legend_to_chart(chart, font_size: Pt):
        """
        Adds a legend with specified font size to a given chart
        """
        chart.has_legend = True
        chart.legend.font.size = Pt(font_size)

    @staticmethod    
    def remove_legend_from_chart(chart):
        """
        Removes the legend from a given chart
        """        
        chart.has_legend = False
        
    def get_all_shapes(self):
        """
        Gets all the shapes from the presentation, including the type of shape for each shape
        Returns a dictionary in format: {slide#: {shape#: shape_type}}
        """
        
        shapes = dict()
        for i, slide in enumerate(self.slides):
            shapes[i] = dict()
            for j, shape in enumerate(slide.shapes):
                shapes[i][j] = shape.shape_type
        return shapes
    
    def get_all_text_from_shapes(self) -> dict:
        """
        Gets all the text from text shapes from the presentation
        Returns a dictionary in format: {slide#: {shape#: shape_text}}
        """        
        shapes = dict()
        for i, slide in enumerate(self.slides):
            shapes[i] = dict()
            for j, shape in enumerate(slide.shapes):
                if shape.has_text_frame:                    
                    shapes[i][j] = shape.text
            if len(shapes[i]) == 0:
                del shapes[i]
        return shapes
    
    def get_all_tables_from_shapes(self) -> dict:
        """
        Gets all the table shapes from the presentation
        Returns a dictionary in format: {slide#: {shape#: 'has_table'}}
        """        
        shapes = dict()
        for i, slide in enumerate(self.slides):
            shapes[i] = dict()
            for j, shape in enumerate(slide.shapes):
                if shape.has_table:                    
                    shapes[i][j] = 'has_table'
            if len(shapes[i]) == 0:
                del shapes[i]
        return shapes
    
    def get_all_charts_from_shapes(self) -> dict:
        """
        Gets all the chart shapes from the presentation
        Returns a dictionary in format: {slide#: {shape#: 'has_chart'}}
        """        
        shapes = dict()
        for i, slide in enumerate(self.slides):
            shapes[i] = dict()
            for j, shape in enumerate(slide.shapes):
                if shape.has_chart:                    
                    shapes[i][j] = 'has_chart'
            if len(shapes[i]) == 0:
                del shapes[i]
        return shapes
    
    def remove_slides_dangerous(self, slide_list: list):
        '''
        Removes slides from the presentation
        Not officially supported and can possibly corrupt the file. Recommended to use it only after all modifications are done
        '''
        slide_list = list(set(slide_list))
        for slide in reversed(sorted(slide_list)):
            slide_ids = list(self._element.sldIdLst)
            self._element.sldIdLst.remove(slide_ids[slide])

    @staticmethod
    def create_series_and_replace_chart_data(shape, categories: list,
                        series_data: dict[list], legend=True, legend_font_size: Union[float, None]=None,
                        chart_color: Union[list[list], None] = None):
        '''
        Replaces chart data with the provided arguments
        Args:
        categories (list): list of the categories for the series
        series_data (dict[list]): dictionary with series label as key and series values as values
        legend (bool): True=display legend, False=no legend
        legend_font_size (float): font size as a float, later converted to Pt
        chart_color (list[list]): list of colors for each series, inside each list 3 integers are required to set the rgb
        '''
        chart_injection_data = CategoryChartData()
        for legend_label, series in series_data.items():
            chart_injection_data.categories(categories)
            chart_injection_data.add_series(legend_label, values=series)            
        shape.chart.replace_data(chart_injection_data)
        shape.chart.has_legend = legend
        if legend_font_size:
            shape.chart.legend.font.size = Pt(legend_font_size)
        if chart_color:
            for i, series in enumerate(shape.chart.series):
                color = RGBColor(*chart_color[i])
                series[i].format.fill.fore_color.rgb = color
                
    def format_table_cell(self, i, j, cell, font_format_methods: dict[dict[dict[list]]]):
        """
        Formats the specified cell using the specified methods if it is specified in font_format_methods
        """
        if j in font_format_methods.keys():
            if i in font_format_methods[j].keys():
                for method_name in font_format_methods[j][i].keys():
                    method = getattr(self, method_name, None)
                    args = [arg for arg in font_format_methods[j][i][method_name] if isinstance(arg, str)]
                    kwargs = [kwarg for kwarg in font_format_methods[j][i][method_name] if isinstance(kwarg, dict)]
                    method(cell, *args, **kwargs)
                    
    def replace_table_data(self, shape, table_data: list[list], font_format_methods: Union[dict[dict[dict[list]]], None] = None):
        """
        Replaces table data with the provided arguments. If the table_data has less rows than the affected table, rows are removed from the table
        to match the length of the new rows.
        ARGS:
        table_data (list[list]): list of rows with list of list of cell values within the rows
        font_format_methods (dict[dict[dict[list]]]): dictionary containing formating methods from this class for each cell
        {row: {column: {method:[*args, {**kwargs]}}}
        """
        while len(shape.table.rows) > len(table_data):
            shape.table._tbl.remove(shape.table.rows[len(shape.table.rows)-1]._tr)
        for j, row in enumerate(shape.table.rows):
            for i, cell in enumerate(row.cells):
                new_text = str(table_data[j][i])
                for run in cell.text_frame.paragraphs[0].runs:
                    run.text = ''
                if len(cell.text_frame.paragraphs[0].runs) == 0:                    
                    run = cell.text_frame.paragraphs[0].add_run()
                else:
                    run = cell.text_frame.paragraphs[0].runs[0]
                run.text = new_text
                if font_format_methods:
                    self.format_table_cell(i, j, cell, font_format_methods)
                
                
                
                